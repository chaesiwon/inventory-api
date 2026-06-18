"""
ppt_exporter.py v6 - ILJIN 임원 보고용 PPT 생성
[v6 수정사항]
 1. 모든 슬라이드의 표는 "행 개수에 맞춰 표 전체 높이를 동적으로 계산"하여
    화면을 넘치거나(overflow) 텅 비는(empty space) 문제를 해결.
 2. 표시되는 모든 금액은 예외 없이 _fmt_amt()를 통과시켜 단위(억원/백만원/원) 일관 적용.
 3. "#" 순번 컬럼 폭을 넉넉히 잡아 두 자리 숫자가 줄바꿈되지 않도록 수정.
 4. 표지/마무리 슬라이드에 일진제강 로고 삽입 (assets/iljin_logo.png).
 5. 행이 많아 한 슬라이드에 다 안 들어가면 자동으로 다음 슬라이드로 이어서 출력.
"""
import io, os, logging
from datetime import datetime
from typing import List, Dict, Any, Optional

logger = logging.getLogger(__name__)

NAVY   = (0x0E, 0x28, 0x41)
BLUE   = (0x15, 0x60, 0x82)
ORANGE = (0xE9, 0x71, 0x32)
RED    = (0xC1, 0x00, 0x1B)
GREEN  = (0x19, 0x6B, 0x24)
WHITE  = (0xFF, 0xFF, 0xFF)
LGRAY  = (0xF4, 0xF6, 0xF9)
MGRAY  = (0xD0, 0xD8, 0xE4)
GRAY   = (0x50, 0x50, 0x50)

LOGO_PATH = os.path.join(os.path.dirname(__file__), "assets", "iljin_logo.png")

SLIDE_W = 13.33
SLIDE_H = 7.5
CONTENT_TOP = 1.68     # 헤더바 아래 본문 시작 y
CONTENT_BOTTOM = 7.1   # 푸터 위 본문 끝 y (이 선을 넘기지 않도록 행 높이를 동적 계산)


def _rgb(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


def _I(n):
    from pptx.util import Inches
    return Inches(n)


def _Pt(n):
    from pptx.util import Pt
    return Pt(n)


def _rect(slide, l, t, w, h, fill, line=None, lw=0.5, rnd=False):
    sh = slide.shapes.add_shape(5 if rnd else 1, _I(l), _I(t), _I(w), _I(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(*fill)
    if line:
        sh.line.color.rgb = _rgb(*line); sh.line.width = _Pt(lw)
    else:
        sh.line.fill.background()
    return sh


def _txt(slide, text, l, t, w, h, sz=12, bold=False, clr=GRAY, align=None, italic=False, valign=None):
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    if align is None:
        align = PP_ALIGN.LEFT
    bx = slide.shapes.add_textbox(_I(l), _I(t), _I(w), _I(h))
    tf = bx.text_frame; tf.word_wrap = True
    if valign is not None:
        tf.vertical_anchor = valign
    lines = str(text).split("\n") if text else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.size = _Pt(sz); run.font.bold = bold
        run.font.color.rgb = _rgb(*clr); run.font.italic = italic
    return bx


def _add_logo(slide, l, t, w, h):
    """일진제강 로고를 비율 유지하며 삽입. 파일이 없으면 조용히 스킵."""
    if not os.path.exists(LOGO_PATH):
        logger.warning(f"로고 파일 없음: {LOGO_PATH}")
        return
    try:
        from PIL import Image
        img = Image.open(LOGO_PATH)
        ratio = img.width / img.height
        # 주어진 박스(w,h) 안에 비율 유지하며 맞춤
        box_ratio = w / h
        if ratio > box_ratio:
            disp_w, disp_h = w, w / ratio
        else:
            disp_w, disp_h = h * ratio, h
        cx = l + (w - disp_w) / 2
        cy = t + (h - disp_h) / 2
        slide.shapes.add_picture(LOGO_PATH, _I(cx), _I(cy), width=_I(disp_w), height=_I(disp_h))
    except Exception as e:
        logger.warning(f"로고 삽입 실패: {e}")


def _fmt_amt(v, unit_label='억원'):
    """[중요] 이 함수는 단위 변환을 하지 않는다. 호출 전에 반드시 api.py의 fmt_amount()를
    거쳐 이미 표시단위(억원/백만원/원)로 변환된 값을 받는다고 가정한다.
    여기서 다시 나누면 이중 변환되어 숫자가 깨지므로 절대 나누기 연산을 추가하지 않는다."""
    try:
        n = float(v or 0)
        if unit_label == '원':
            return f"{round(n):,}원"
        return f"{n:,.2f}{unit_label}"
    except Exception:
        return "-"


def _fmt_wt(v):
    try:
        return f"{float(v or 0):.1f}ton"
    except Exception:
        return "-"


def _fmt_cnt(v):
    try:
        return f"{int(v or 0):,}건"
    except Exception:
        return "-"


def _fmt_pct(v):
    try:
        return f"{float(v or 0):.1f}%"
    except Exception:
        return "-"


def _hbar(slide, title, sub=None):
    _rect(slide, 0, 0, SLIDE_W, 1.52, NAVY)
    _rect(slide, 0, 1.48, SLIDE_W, 0.06, ORANGE)
    _txt(slide, title, 0.48, 0.2, 11, 0.78, sz=24, bold=True, clr=WHITE)
    if sub:
        _txt(slide, sub, 0.48, 0.9, 11, 0.42, sz=11, clr=MGRAY, italic=True)


def _foot(slide, n=None, label="장기재고 소진계획 계획/실적 비교 보고서"):
    from pptx.enum.text import PP_ALIGN
    _rect(slide, 0, 7.2, SLIDE_W, 0.3, NAVY)
    _txt(slide, label, 0.42, 7.22, 9, 0.26, sz=9, clr=MGRAY)
    if n:
        _txt(slide, str(n), 12.55, 7.22, 0.65, 0.26, sz=9, clr=MGRAY, align=PP_ALIGN.RIGHT)


def _draw_table(slide, headers, rows, col_widths, x0=0.38, y0=None,
                 header_color=NAVY, max_height=None, font_size=10, header_font_size=11):
    """
    행 개수에 맞춰 표 전체 높이를 동적으로 계산해서 그리는 공통 테이블 함수.
    - 사용 가능한 세로 공간(max_height) 안에서 행 높이를 균등 분배(최소/최대 제한 적용).
    - 이렇게 하면 행이 적을 때 화면 아래가 텅 비거나, 행이 많을 때 화면을 넘치는 문제를 방지.
    rows: List[List[str]] (이미 포맷된 문자열)
    """
    from pptx.enum.text import PP_ALIGN

    if y0 is None:
        y0 = CONTENT_TOP
    if max_height is None:
        max_height = CONTENT_BOTTOM - y0

    n_rows = len(rows)
    header_h = 0.46
    if n_rows == 0:
        avail = max_height - header_h
        row_h = avail
    else:
        avail = max_height - header_h
        row_h = avail / n_rows
        # 행 높이가 너무 크면(행이 적을 때) 적당한 상한을 둬서 카드처럼 늘어지지 않게 함
        row_h = min(row_h, 0.62)
        row_h = max(row_h, 0.34)

    xs = [x0]
    for w in col_widths[:-1]:
        xs.append(xs[-1] + w)

    # 헤더
    for h, x, w in zip(headers, xs, col_widths):
        _rect(slide, x, y0, w, header_h, header_color)
        _txt(slide, h, x + 0.05, y0 + (header_h - 0.3) / 2, w - 0.08, 0.3,
             sz=header_font_size, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)

    # 데이터 행
    for ri, row in enumerate(rows):
        y = y0 + header_h + ri * row_h
        bg = LGRAY if ri % 2 == 0 else WHITE
        for ci, (val, x, w) in enumerate(zip(row, xs, col_widths)):
            _rect(slide, x, y, w, row_h, bg, line=MGRAY, lw=0.4)
            cell = val if isinstance(val, str) else val.get("text", "")
            clr = val.get("color", GRAY) if isinstance(val, dict) else GRAY
            bold = val.get("bold", False) if isinstance(val, dict) else False
            align = val.get("align", PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT) if isinstance(val, dict) else (PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)
            fs = min(font_size, 11) if row_h < 0.4 else font_size
            _txt(slide, cell, x + 0.06, y + max(0.02, (row_h - fs / 72 * 1.3) / 2), w - 0.1, row_h - 0.04,
                 sz=fs, clr=clr, bold=bold, align=align)

    return y0 + header_h + n_rows * row_h  # 표가 끝나는 y좌표 반환


def _paginate(rows, max_rows_per_slide):
    """행이 많으면 슬라이드 여러 장으로 나눈다."""
    if not rows:
        return [[]]
    return [rows[i:i + max_rows_per_slide] for i in range(0, len(rows), max_rows_per_slide)]


def _slide_cover(prs, ref_date, generated_at, title="장기재고 소진계획", subtitle="계획/실적 비교 보고서"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _rect(s, 9.85, 0, 3.48, SLIDE_H, BLUE)
    _rect(s, 9.8, 0, 0.07, SLIDE_H, ORANGE)

    # 로고 (우측 패널 상단, 흰 배경 카드 위에)
    _rect(s, 10.2, 0.55, 2.7, 0.95, WHITE, rnd=True)
    _add_logo(s, 10.35, 0.68, 2.4, 0.7)

    _txt(s, title, 0.65, 2.15, 9, 0.88, sz=40, bold=True, clr=WHITE)
    _txt(s, subtitle, 0.65, 3.0, 9, 0.88, sz=32, bold=True, clr=WHITE)
    rd = f"{ref_date[:4]}-{ref_date[4:6]}-{ref_date[6:8]}" if len(ref_date) == 8 else ref_date
    _rect(s, 0.65, 3.95, 5.5, 0.07, ORANGE)
    _txt(s, f"기준일: {rd}", 0.65, 4.12, 9, 0.48, sz=15, clr=(0xCC, 0xDD, 0xFF))
    _txt(s, f"생성: {generated_at}", 0.65, 5.0, 9, 0.38, sz=12, clr=MGRAY)
    _txt(s, "일진제강", 10.2, 6.7, 2.7, 0.4, sz=12, bold=True, clr=(0xCC, 0xDD, 0xFF))


def _slide_kpi(prs, summary, ref_date, unit_label):
    from pptx.enum.text import PP_ALIGN
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _hbar(s, "01  핵심 KPI 요약")
    pt = summary.get("plan_total", 0) or 1
    ac = summary.get("action_count", 0) or 0
    nc = summary.get("no_action_count", 0) or 0
    rate_cnt = summary.get("action_rate", 0)

    kpis = [
        ("계획 등록 LOT", _fmt_cnt(pt), _fmt_amt(summary.get("total_amount"), unit_label), NAVY),
        ("조치 완료", _fmt_cnt(ac), _fmt_amt(summary.get("action_amount"), unit_label), GREEN),
        ("미 조치", _fmt_cnt(nc), _fmt_amt(summary.get("no_action_amount"), unit_label), RED),
        ("달성률", _fmt_pct(rate_cnt), "건수 기준", BLUE),
        ("소진완료금액", _fmt_amt(summary.get("consumed_amount"), unit_label), "LOT단가×실적중량", ORANGE),
    ]
    card_w = 2.35
    gap = (SLIDE_W - 0.38 * 2 - card_w * 5) / 4
    for i, (lbl, val, sub, col) in enumerate(kpis):
        lx = 0.38 + i * (card_w + gap)
        _rect(s, lx, 1.85, card_w, 1.55, col, rnd=True)
        _txt(s, lbl, lx + 0.1, 1.95, card_w - 0.2, 0.45, sz=11, bold=True, clr=(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
        _txt(s, val, lx + 0.1, 2.42, card_w - 0.2, 0.55, sz=19, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, sub, lx + 0.1, 3.0, card_w - 0.2, 0.3, sz=9, clr=(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

    bar_y = 3.85
    _txt(s, f"달성률 (조치 {ac:,}건 / 전체 {pt:,}건)", 0.4, bar_y - 0.42, 8, 0.35, sz=13, bold=True, clr=GRAY)
    _rect(s, 0.38, bar_y, 12.57, 0.85, LGRAY, rnd=True)
    fill_w = max(0.3, 12.57 * ac / pt)
    fill_c = GREEN if rate_cnt >= 70 else (0xF5, 0x9E, 0x0B) if rate_cnt >= 40 else RED
    _rect(s, 0.38, bar_y, fill_w, 0.85, fill_c, rnd=True)
    # 달성률 숫자는 막대 폭과 무관하게 막대 우측 바깥에 고정 배치하여 항상 잘 보이게 함
    pct_color = fill_c
    _txt(s, _fmt_pct(rate_cnt), 11.3, bar_y + 0.18, 1.6, 0.5, sz=22, bold=True, clr=pct_color, align=None)

    # 하단 여백을 활용한 추가 안내 (미조치 현황 요약으로 빈 공간 방지)
    info_y = 5.1
    _rect(s, 0.38, info_y, 12.57, 1.5, LGRAY, rnd=True)
    _txt(s, "보고서 안내", 0.62, info_y + 0.18, 6, 0.35, sz=13, bold=True, clr=NAVY)
    _txt(s,
         "· 계획은 시스템에 등록된 소진계획(담당부서/사유/방안/기한)을 기준으로 합니다.\n"
         "· 실적은 업로드된 장기재고현황 파일의 상세시트(재고_상세/재공_상세)에서 조회기준일별로 산출됩니다.\n"
         "· 계획과 실적은 LOT NO를 기준으로 1:1 매칭하여 비교합니다.",
         0.62, info_y + 0.55, 12.0, 0.9, sz=11, clr=GRAY)
    _foot(s, 2)


def _slide_type_compare(prs, summary, unit_label, page_num_start=3):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _hbar(s, "02  유형별 계획 vs 실적 비교")
    plan_by = summary.get("plan_by_type", [])
    actual_by = summary.get("actual_by_type", [])
    pM = {(r.get("plan_type") or "미등록"): r for r in plan_by}
    aM = {(r.get("actual_type") or "기타"): r for r in actual_by}
    all_types = sorted(set(pM.keys()) | set(aM.keys()))

    headers = ["구분", "계획 건수", "계획 중량(ton)", "실적 건수", "실적 중량(ton)", "달성률(중량)"]
    col_widths = [2.7, 1.9, 2.1, 1.9, 2.1, 1.87]

    rows = []
    for typ in all_types:
        pm = pM.get(typ, {}); am = aM.get(typ, {})
        pc = int(pm.get("plan_count", 0)); pw = float(pm.get("plan_weight", 0))
        ac = int(am.get("actual_count", 0)); aw = float(am.get("actual_weight", 0))
        pct = round(aw / pw * 100, 1) if pw > 0 else 0
        pct_color = GREEN if pct >= 70 else RED if pct < 40 else GRAY
        rows.append([
            {"text": typ, "align": None},
            f"{pc:,}", f"{pw:.1f}", f"{ac:,}", f"{aw:.1f}",
            {"text": f"{pct:.1f}%", "color": pct_color, "bold": True},
        ])

    _draw_table(s, headers, rows, col_widths)
    _foot(s, page_num_start)


def _slide_top_items(prs, items, unit_label, page_num_start=4):
    """행이 많으면 여러 슬라이드로 분할."""
    pages = _paginate(items, max_rows_per_slide=12)
    slides_made = []
    headers = ["#", "공장", "LOT NO", "품명", f"금액({unit_label})", "계획유형", "조치여부"]
    col_widths = [0.5, 1.05, 1.85, 4.15, 1.5, 1.45, 1.42]

    for pi, page_items in enumerate(pages):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        title = "03  고액 상위 장기재고 현황" + (f" ({pi+1}/{len(pages)})" if len(pages) > 1 else "")
        _hbar(s, title)
        if not page_items:
            _txt(s, "데이터가 없습니다.", 4, 4, 5, 0.8, sz=18, clr=GRAY)
            _foot(s, page_num_start + pi)
            slides_made.append(s)
            continue
        rows = []
        for idx, item in enumerate(page_items):
            global_idx = pi * 12 + idx + 1
            amt = _fmt_amt(item.get("amount", 0), unit_label)
            act = item.get("action_status", "")
            act_color = GREEN if act == "조치" else RED
            rows.append([
                str(global_idx),
                str(item.get("factory", ""))[:8],
                str(item.get("lot_no", ""))[:16],
                {"text": str(item.get("item_name", ""))[:28], "align": None},
                amt,
                str(item.get("plan_type", "") or "-")[:10],
                {"text": act, "color": act_color, "bold": True},
            ])
        _draw_table(s, headers, rows, col_widths)
        _foot(s, page_num_start + pi)
        slides_made.append(s)
    return len(pages)


def _slide_plan_trend(prs, plan_trend, unit_label, page_num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _hbar(s, "04  월별 소진계획 현황")
    if not plan_trend:
        _txt(s, "소진계획 데이터가 없습니다.", 4, 4, 5, 0.8, sz=16, clr=GRAY)
        _foot(s, page_num); return

    headers = ["소진계획월", "계획 건수", "계획 중량(ton)", f"계획 금액({unit_label})"]
    col_widths = [2.6, 2.0, 2.5, 2.47]  # 합계 9.57 + x0(1.0)*2 여백 = 11.57, 슬라이드폭(13.33) 이내
    rows = []
    for row in plan_trend[:12]:
        rows.append([
            str(row.get("plan_month", ""))[:7],
            f"{int(row.get('plan_count', 0)):,}",
            f"{float(row.get('plan_weight_ton', 0)):.1f}",
            _fmt_amt(row.get("plan_amount", 0), unit_label),
        ])
    _draw_table(s, headers, rows, col_widths, x0=1.5)
    _foot(s, page_num)


def _slide_cost_center(prs, cc_items, unit_label, page_num):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _hbar(s, "05  원가중심점별 장기재고 현황")
    if not cc_items:
        _txt(s, "데이터가 없습니다.", 4, 4, 5, 0.8, sz=16, clr=GRAY)
        _foot(s, page_num); return

    headers = ["원가중심점", "건수", "중량(ton)", f"금액({unit_label})", "계획등록", "실적확인", "달성률"]
    col_widths = [2.7, 1.0, 1.5, 1.8, 1.3, 1.3, 1.57]
    rows = []
    for item in cc_items[:12]:
        ic = int(item.get("item_count", 0)); pl = int(item.get("plan_count", 0)); ac = int(item.get("actual_count", 0))
        rate = round(ac / ic * 100, 1) if ic > 0 else 0
        rate_color = GREEN if rate >= 70 else RED if rate < 40 else GRAY
        rows.append([
            {"text": str(item.get("cc_name", ""))[:20], "align": None},
            f"{ic:,}", f"{float(item.get('total_weight', 0)):.1f}",
            _fmt_amt(item.get("total_amount", 0), unit_label),
            f"{pl:,}", f"{ac:,}",
            {"text": f"{rate:.1f}%", "color": rate_color, "bold": True},
        ])
    _draw_table(s, headers, rows, col_widths)
    _foot(s, page_num)


def _slide_closing(prs, generated_at):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    _rect(s, 9.85, 0, 3.48, SLIDE_H, BLUE)
    _rect(s, 9.8, 0, 0.07, SLIDE_H, ORANGE)
    _rect(s, 10.2, 0.55, 2.7, 0.95, WHITE, rnd=True)
    _add_logo(s, 10.35, 0.68, 2.4, 0.7)
    _txt(s, "감사합니다", 0.65, 2.5, 9, 1.05, sz=52, bold=True, clr=WHITE)
    _txt(s, f"생성일시: {generated_at}", 0.65, 4.0, 6, 0.35, sz=11, clr=(0x60, 0x70, 0x80))
    _txt(s, "일진제강", 10.2, 6.7, 2.7, 0.4, sz=12, bold=True, clr=(0xCC, 0xDD, 0xFF))


INV_FOOT_LABEL = "장기재고현황 조회 보고서"


def _slide_inventory_list(prs, items, unit_label, page_num_start, title_prefix="03  장기재고현황 상세 목록"):
    """장기재고현황 조회 결과를 표로 출력. 행이 많으면 슬라이드 분할."""
    pages = _paginate(items, max_rows_per_slide=12)
    headers = ["#", "공장", "LOT NO", "품명", f"금액({unit_label})", "개월", "계획유형"]
    col_widths = [0.5, 1.05, 1.85, 4.15, 1.5, 1.45, 1.42]

    for pi, page_items in enumerate(pages):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        title = title_prefix + (f" ({pi+1}/{len(pages)})" if len(pages) > 1 else "")
        _hbar(s, title)
        if not page_items:
            _txt(s, "데이터가 없습니다.", 4, 4, 5, 0.8, sz=18, clr=GRAY)
            _foot(s, page_num_start + pi, label=INV_FOOT_LABEL)
            continue
        rows = []
        for idx, item in enumerate(page_items):
            global_idx = pi * 12 + idx + 1
            amt = _fmt_amt(item.get("amount", 0), unit_label)
            months = item.get("months_label") or "-"
            is_critical = months == "7개월이상"
            rows.append([
                str(global_idx),
                str(item.get("factory", ""))[:8],
                str(item.get("lot_no", ""))[:16],
                {"text": str(item.get("item_name", ""))[:28], "align": None},
                amt,
                {"text": months, "color": RED if is_critical else GRAY, "bold": is_critical},
                str(item.get("plan_type", "") or "-")[:10],
            ])
        _draw_table(s, headers, rows, col_widths)
        _foot(s, page_num_start + pi, label=INV_FOOT_LABEL)
    return len(pages)


def _slide_inventory_summary(prs, summary, ref_date, unit_label):
    from pptx.enum.text import PP_ALIGN
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _hbar(s, "01  장기재고현황 요약")

    kpis = [
        ("총 장기재고", _fmt_cnt(summary.get("total_count")), _fmt_amt(summary.get("total_amount"), unit_label), NAVY),
        ("총 중량", _fmt_wt(summary.get("total_weight")), "전체 합계", BLUE),
        ("7개월이상", _fmt_cnt(summary.get("critical_count")), _fmt_amt(summary.get("critical_amount"), unit_label), RED),
        ("계획 등록", _fmt_cnt(summary.get("plan_count")), "건", GREEN),
    ]
    card_w = 2.8
    gap = (SLIDE_W - 0.38 * 2 - card_w * 4) / 3
    for i, (lbl, val, sub, col) in enumerate(kpis):
        lx = 0.38 + i * (card_w + gap)
        _rect(s, lx, 1.85, card_w, 1.55, col, rnd=True)
        _txt(s, lbl, lx + 0.1, 1.95, card_w - 0.2, 0.45, sz=12, bold=True, clr=(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
        _txt(s, val, lx + 0.1, 2.42, card_w - 0.2, 0.55, sz=20, bold=True, clr=WHITE, align=PP_ALIGN.CENTER)
        _txt(s, sub, lx + 0.1, 3.0, card_w - 0.2, 0.3, sz=10, clr=(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

    info_y = 3.85
    _rect(s, 0.38, info_y, 12.57, 1.5, LGRAY, rnd=True)
    _txt(s, "보고서 안내", 0.62, info_y + 0.18, 6, 0.35, sz=13, bold=True, clr=NAVY)
    rd = f"{ref_date[:4]}-{ref_date[4:6]}-{ref_date[6:8]}" if len(str(ref_date)) == 8 else ref_date
    _txt(s,
         f"· 조회기준일: {rd} 기준으로 집계된 장기재고현황입니다.\n"
         "· 저장품(item_type='저장품')은 장기재고관리 대상에서 제외되어 본 보고서에 포함되지 않습니다.\n"
         "· 7개월이상 장기화된 재고는 빨간색으로 강조 표시되어 우선적인 소진계획 검토가 필요합니다.",
         0.62, info_y + 0.55, 12.0, 0.9, sz=11, clr=GRAY)
    _foot(s, 2, label=INV_FOOT_LABEL)


def generate_inventory_ppt(
    items: List[Dict[str, Any]],
    summary: dict,
    ref_date: str = "",
    unit_label: str = "억원",
) -> bytes:
    """장기재고현황 조회 화면용 PPT 생성 (요구사항 7)"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise ImportError("python-pptx 설치 필요")

    generated_at = datetime.now().strftime("%Y-%m-%d %H:%M")
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    _slide_cover(prs, ref_date, generated_at, title="장기재고현황", subtitle="조회 보고서")
    _slide_inventory_summary(prs, summary, ref_date, unit_label)
    n_pages = _slide_inventory_list(prs, items, unit_label, page_num_start=3)
    _slide_closing(prs, generated_at)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    result = buf.read()
    logger.info(f"재고현황 PPT 생성 완료: {len(result):,} bytes, {n_pages}개 목록 슬라이드")
    return result


def generate_compare_ppt(
    summary: dict,
    items: List[Dict[str, Any]],
    ref_date: str = "",
    plan_trend: Optional[List] = None,
    cc_items: Optional[List] = None,
    unit_label: str = "억원",
) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise ImportError("python-pptx 설치 필요")

    generated_at = datetime.now().strftime("%Y-%m-%d %H:%M")
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    _slide_cover(prs, ref_date, generated_at)
    _slide_kpi(prs, summary, ref_date, unit_label)
    _slide_type_compare(prs, summary, unit_label, page_num_start=3)
    n_top_pages = _slide_top_items(prs, items, unit_label, page_num_start=4)
    next_page = 4 + n_top_pages
    _slide_plan_trend(prs, plan_trend or [], unit_label, page_num=next_page)
    _slide_cost_center(prs, cc_items or [], unit_label, page_num=next_page + 1)
    _slide_closing(prs, generated_at)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    result = buf.read()
    logger.info(f"PPT 생성 완료: {len(result):,} bytes")
    return result
